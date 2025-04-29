import os
from pptx import Presentation
import json
import re
from pathlib import Path

def extract_pptx_content(pptx_path):
    """Extract content from PowerPoint file"""
    prs = Presentation(pptx_path)
    
    slides_content = []
    
    for i, slide in enumerate(prs.slides):
        slide_content = {
            "slide_number": i + 1,
            "title": "",
            "content": [],
            "notes": "",
            "images": []
        }
        
        # Extract title
        if slide.shapes.title:
            slide_content["title"] = slide.shapes.title.text
        
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
                # Skip empty or whitespace-only text
                if shape.text.strip():
                    slide_content["content"].append(shape.text)
        
        # Extract notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            for shape in notes_slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_content["notes"] += shape.text
        
        # Track image shapes for later extraction
        for j, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image_filename = f"slide_{i+1}_image_{j+1}.png"
                slide_content["images"].append(image_filename)
        
        slides_content.append(slide_content)
    
    return slides_content

def save_images(pptx_path, output_dir):
    """Extract and save images from PowerPoint file"""
    prs = Presentation(pptx_path)
    
    # Create images directory if it doesn't exist
    images_dir = os.path.join(output_dir, "images")
    os.makedirs(images_dir, exist_ok=True)
    
    image_paths = []
    
    for i, slide in enumerate(prs.slides):
        for j, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image = shape.image
                image_bytes = image.blob
                
                # Save image
                image_filename = f"slide_{i+1}_image_{j+1}.png"
                image_path = os.path.join(images_dir, image_filename)
                
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
                
                image_paths.append(image_path)
    
    return image_paths

def main():
    pptx_path = "/home/ubuntu/upload/1_sissejuhatus.pptx"
    output_dir = "/home/ubuntu/website_project"
    
    # Extract content
    slides_content = extract_pptx_content(pptx_path)
    
    # Save content to JSON file
    with open(os.path.join(output_dir, "pptx_content.json"), "w", encoding="utf-8") as f:
        json.dump(slides_content, f, ensure_ascii=False, indent=2)
    
    # Extract and save images
    image_paths = save_images(pptx_path, output_dir)
    
    print(f"Extracted content from {len(slides_content)} slides")
    print(f"Saved {len(image_paths)} images")

if __name__ == "__main__":
    main()
